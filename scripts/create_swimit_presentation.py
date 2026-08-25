from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PUBLIC = ROOT / "client" / "public"
CURSOR_ASSETS = Path(
    r"C:\Users\sneha\.cursor\projects\d-projects-SwimIT-swimIT\assets"
)
OUTPUT = ROOT / "docs" / "SwimIT_Digital_Transformation_Presentation.pptx"

NAVY = RGBColor(7, 45, 92)
BLUE = RGBColor(27, 112, 214)
CYAN = RGBColor(0, 174, 224)
GREEN = RGBColor(18, 151, 91)
ORANGE = RGBColor(242, 128, 39)
RED = RGBColor(209, 67, 67)
INK = RGBColor(30, 49, 70)
MUTED = RGBColor(91, 111, 132)
PALE = RGBColor(238, 247, 255)
PALE_GREEN = RGBColor(235, 248, 241)
PALE_ORANGE = RGBColor(255, 246, 233)
LIGHT = RGBColor(248, 251, 254)
LINE = RGBColor(211, 224, 236)
WHITE = RGBColor(255, 255, 255)


def add_rect(slide, x, y, w, h, fill, radius=False, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    font="Aptos",
    valign=MSO_ANCHOR.TOP,
    margin=0.04,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, size=17, color=INK, spacing=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
        p.line_spacing = 1.05
        p.text = "•  " + item
    return box


def image_cover(slide, path, x, y, w, h):
    from PIL import Image

    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    image_ratio = iw / ih
    box_ratio = w / h
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if image_ratio > box_ratio:
        visible = box_ratio / image_ratio
        pic.crop_left = (1 - visible) / 2
        pic.crop_right = (1 - visible) / 2
    else:
        visible = image_ratio / box_ratio
        pic.crop_top = (1 - visible) / 2
        pic.crop_bottom = (1 - visible) / 2
    return pic


def add_title(slide, title, subtitle=None, number=None):
    add_text(slide, title, 0.55, 0.28, 11.9, 0.55, 25, NAVY, True)
    add_rect(slide, 0.55, 0.91, 1.15, 0.05, CYAN)
    if subtitle:
        add_text(slide, subtitle, 0.55, 1.0, 12.0, 0.4, 11.5, MUTED)
    if number is not None:
        add_text(slide, f"{number:02}", 12.35, 0.32, 0.38, 0.3, 10, MUTED, True, PP_ALIGN.RIGHT)


def add_footer(slide, number):
    add_rect(slide, 0.55, 7.18, 12.2, 0.012, LINE)
    add_text(slide, "SwimIT · Swimming Pool Management System", 0.55, 7.22, 5.4, 0.22, 8.5, MUTED)
    add_text(slide, str(number), 12.2, 7.22, 0.55, 0.22, 8.5, MUTED, align=PP_ALIGN.RIGHT)


def add_metric(slide, value, label, x, y, w, color=BLUE):
    add_rect(slide, x, y, w, 1.04, WHITE, True, LINE)
    add_text(slide, value, x + 0.18, y + 0.13, w - 0.36, 0.38, 24, color, True)
    add_text(slide, label, x + 0.18, y + 0.58, w - 0.36, 0.27, 10.5, MUTED)


def add_icon_card(slide, title, body, x, y, w, h, color=BLUE, glyph="1"):
    add_rect(slide, x, y, w, h, WHITE, True, LINE)
    add_rect(slide, x + 0.18, y + 0.18, 0.48, 0.48, color, True)
    add_text(
        slide,
        glyph,
        x + 0.18,
        y + 0.18,
        0.48,
        0.48,
        13,
        WHITE,
        True,
        PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(slide, title, x + 0.78, y + 0.17, w - 0.94, 0.34, 14.5, NAVY, True)
    add_text(slide, body, x + 0.18, y + 0.77, w - 0.36, h - 0.92, 11.5, MUTED)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# 1 — Cover
slide = prs.slides.add_slide(blank)
image_cover(slide, PUBLIC / "marketing-hero-swimmer.jpg", 7.2, 0, 6.13, 7.5)
add_rect(slide, 0, 0, 7.65, 7.5, WHITE)
slide.shapes.add_picture(str(PUBLIC / "swimit-wordmark.png"), Inches(0.72), Inches(0.55), width=Inches(3.75))
add_text(slide, "From paper registers\nto one smart pool system", 0.72, 2.0, 6.35, 1.55, 30, NAVY, True)
add_text(
    slide,
    "Digital registration, pass, invoice, attendance and accounts — built for swimming pools.",
    0.74,
    3.75,
    5.9,
    0.8,
    17,
    MUTED,
)
add_rect(slide, 0.74, 4.87, 3.45, 0.58, BLUE, True)
add_text(
    slide,
    "One login · One record · Better control",
    0.88,
    5.0,
    3.2,
    0.27,
    12,
    WHITE,
    True,
    PP_ALIGN.CENTER,
)
add_text(slide, "Presentation for swimming pool owners and operators", 0.74, 6.55, 5.8, 0.28, 10.5, MUTED)


# 2 — Current reality
slide = prs.slides.add_slide(blank)
add_title(slide, "The current paper-based reality", "The same information is written repeatedly, but visibility remains low.", 2)
items = [
    ("Registration form", "Details stay in files; searching takes time.", "F"),
    ("Physical register", "Manual gate entry and attendance marking.", "R"),
    ("Bill book", "Duplicate writing, carbon copies and month-end totals.", "B"),
    ("Printed pass", "Can be lost, forgotten, copied or used after expiry.", "P"),
    ("Expense register", "Profit is known only after manual calculation.", "₹"),
    ("Scattered records", "Owner depends on staff and multiple notebooks.", "!"),
]
for idx, (title, body, glyph) in enumerate(items):
    col, row = idx % 3, idx // 3
    add_icon_card(slide, title, body, 0.62 + col * 4.17, 1.55 + row * 2.25, 3.82, 1.82, [BLUE, ORANGE, RED][col], glyph)
add_rect(slide, 0.62, 6.18, 12.08, 0.62, PALE, True, PALE)
add_text(
    slide,
    "The result: more writing, slower service, payment disputes and no real-time view of the pool.",
    0.88,
    6.34,
    11.55,
    0.26,
    14,
    NAVY,
    True,
    PP_ALIGN.CENTER,
)
add_footer(slide, 2)


# 3 — Direct replacement
slide = prs.slides.add_slide(blank)
add_title(slide, "Paper process → SwimIT", "A direct replacement for the registers pools already use.", 3)
rows = [
    ("Paper registration form", "Digital registration form with photo, contact and batch"),
    ("Swimmer register", "Searchable Active / Inactive swimmer list"),
    ("Bill book", "Payment record and tax-inclusive invoice"),
    ("Printed pass", "Digital QR pass sent on WhatsApp"),
    ("Gate attendance register", "Pass scanner and automatic attendance"),
    ("Expense notebook", "Pool Expenses and Balance Sheet"),
]
for i, (before, after) in enumerate(rows):
    y = 1.45 + i * 0.82
    add_rect(slide, 0.65, y, 4.65, 0.62, LIGHT if i % 2 == 0 else WHITE, True, LINE)
    add_text(slide, before, 0.92, y + 0.15, 4.1, 0.27, 13, INK, i == 0)
    add_text(slide, "→", 5.52, y + 0.13, 0.55, 0.3, 19, CYAN, True, PP_ALIGN.CENTER)
    add_rect(slide, 6.3, y, 6.4, 0.62, PALE_GREEN if i % 2 == 0 else WHITE, True, LINE)
    add_text(slide, after, 6.57, y + 0.15, 5.88, 0.27, 13, NAVY, True)
add_footer(slide, 3)


# 4 — Workflow image
slide = prs.slides.add_slide(blank)
add_title(slide, "One connected workflow", "Information is entered once and reused across every operation.", 4)
slide.shapes.add_picture(str(PUBLIC / "swimit-workflow.png"), Inches(0.58), Inches(1.3), width=Inches(12.15))
add_rect(slide, 1.1, 6.55, 11.15, 0.42, PALE, True, PALE)
add_text(
    slide,
    "Register → Take payment → Send pass & invoice → Scan at gate → Review attendance and accounts",
    1.25,
    6.65,
    10.85,
    0.22,
    11.5,
    NAVY,
    True,
    PP_ALIGN.CENTER,
)
add_footer(slide, 4)


# 5 — Owner benefits
slide = prs.slides.add_slide(blank)
add_title(slide, "What the pool owner gains", "Control the operation without waiting for registers to be totalled.", 5)
benefits = [
    ("Daily visibility", "Active swimmers, attendance, expiring passes and payment totals.", "1"),
    ("Less revenue leakage", "Expired or unpaid passes are visible before entry.", "2"),
    ("Financial control", "Pass income, expenses and coach payouts in one view.", "3"),
    ("Accountability", "Separate staff logins, page access and activity history.", "4"),
    ("Professional image", "Pool branding on digital pass and invoice.", "5"),
    ("Scalable operation", "Growth adds records — not more books and manual totals.", "6"),
]
for idx, (title, body, glyph) in enumerate(benefits):
    col, row = idx % 2, idx // 2
    add_icon_card(slide, title, body, 0.72 + col * 6.18, 1.48 + row * 1.68, 5.73, 1.36, BLUE if col == 0 else GREEN, glyph)
add_footer(slide, 5)


# 6 — Staff benefits
slide = prs.slides.add_slide(blank)
add_title(slide, "Faster work for desk, gate and coaches", "Every role gets the information needed for the task.", 6)
roles = [
    ("FRONT DESK", BLUE, ["Register once", "Cash or UPI payment", "Issue / resend pass and invoice", "Search by name or mobile"]),
    ("GATE STAFF", GREEN, ["Scan QR / ID", "Confirm pass validity", "Optional photo check", "Attendance updates automatically"]),
    ("COACH / MANAGER", ORANGE, ["Batch and swimmer lists", "Monthly attendance", "Coach payment summary", "Operational reports"]),
]
for idx, (role, color, bullets) in enumerate(roles):
    x = 0.62 + idx * 4.18
    add_rect(slide, x, 1.5, 3.78, 4.75, WHITE, True, LINE)
    add_rect(slide, x, 1.5, 3.78, 0.7, color, True, color)
    add_text(slide, role, x + 0.25, 1.7, 3.28, 0.28, 14, WHITE, True, PP_ALIGN.CENTER)
    add_bullets(slide, bullets, x + 0.3, 2.53, 3.2, 2.75, 15, INK, 15)
add_rect(slide, 1.75, 6.45, 9.82, 0.45, PALE, True, PALE)
add_text(slide, "Less writing · Fewer errors · Shorter queues · Clear responsibility", 1.9, 6.56, 9.52, 0.22, 12.5, NAVY, True, PP_ALIGN.CENTER)
add_footer(slide, 6)


# 7 — Swimmer experience with screenshots
slide = prs.slides.add_slide(blank)
add_title(slide, "A better experience for swimmers and parents", "From payment to active pass without repeated visits to the desk.", 7)
shot_active = CURSOR_ASSETS / "c__Users_sneha_AppData_Roaming_Cursor_User_workspaceStorage_26803b28d5ecc080f6c90797d58fa982_images_image-1f6a37ce-2a5b-4f61-aa89-a736857bfcb8.png"
image_cover(slide, shot_active, 5.9, 1.45, 6.75, 3.45)
add_rect(slide, 0.68, 1.45, 4.65, 3.45, PALE, True, PALE)
add_text(slide, "Swimmer journey", 1.0, 1.78, 3.95, 0.35, 18, NAVY, True)
add_bullets(
    slide,
    [
        "Fill the registration form on a phone",
        "Pay by cash or UPI",
        "Receive pass and invoice on WhatsApp",
        "Show the QR pass at the gate",
        "Receive reminders before expiry",
    ],
    0.95,
    2.35,
    4.05,
    2.05,
    14.5,
)
add_text(slide, "Actual SwimIT Active Swimmer view", 7.8, 4.98, 3.0, 0.25, 9.5, MUTED, False, PP_ALIGN.CENTER)
add_metric(slide, "No paper to carry", "Digital pass stays on the swimmer’s phone", 0.7, 5.45, 3.75, BLUE)
add_metric(slide, "Clear validity", "Expiry date visible to swimmer and gate", 4.78, 5.45, 3.75, GREEN)
add_metric(slide, "Easy support", "Pass and invoice can be resent", 8.86, 5.45, 3.75, ORANGE)
add_footer(slide, 7)


# 8 — Modules
slide = prs.slides.add_slide(blank)
add_title(slide, "Everything needed to run the pool", "Modules are organised around setup, daily operations, information and finance.", 8)
modules = [
    ("SETUP", BLUE, "Core Info\nBatches & timings\nPass types\nHolidays\nUser access"),
    ("DAILY OPERATIONS", GREEN, "Registration\nPass Payment\nPass Scanner\nWhatsApp\nWater Quality"),
    ("INFORMATION", CYAN, "Dashboard\nSwimmer’s List\nAttendance Sheet\nStaff List\nPayment Details"),
    ("FINANCE", ORANGE, "Pool Expenses\nCoach Payment\nInvoices\nBalance Sheet\nDownloads"),
]
for idx, (head, color, body) in enumerate(modules):
    x = 0.55 + idx * 3.17
    add_rect(slide, x, 1.55, 2.87, 4.85, WHITE, True, LINE)
    add_rect(slide, x, 1.55, 2.87, 0.72, color, True, color)
    add_text(slide, head, x + 0.12, 1.77, 2.63, 0.26, 12.5, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, body, x + 0.35, 2.65, 2.17, 2.95, 15, INK, False, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
add_footer(slide, 8)


# 9 — Finance screenshot
slide = prs.slides.add_slide(blank)
add_title(slide, "Money visibility without manual totals", "Every pass payment becomes a credit; expenses and payouts become debits.", 9)
shot_balance = CURSOR_ASSETS / "c__Users_sneha_AppData_Roaming_Cursor_User_workspaceStorage_26803b28d5ecc080f6c90797d58fa982_images_image-a3a7b4aa-b06a-476d-ae6a-585eb4981959.png"
image_cover(slide, shot_balance, 0.65, 1.38, 8.35, 4.45)
add_text(slide, "Actual SwimIT Balance Sheet", 3.15, 5.91, 3.2, 0.24, 9.5, MUTED, False, PP_ALIGN.CENTER)
add_rect(slide, 9.35, 1.38, 3.3, 4.45, PALE_GREEN, True, PALE_GREEN)
add_text(slide, "Owner can answer instantly:", 9.72, 1.72, 2.6, 0.55, 17, NAVY, True)
add_bullets(
    slide,
    [
        "How much was collected?",
        "What was spent?",
        "Which swimmer paid?",
        "Cash or online?",
        "What is the closing balance?",
    ],
    9.62,
    2.55,
    2.65,
    2.7,
    14.5,
    INK,
    13,
)
add_rect(slide, 9.62, 5.25, 2.7, 0.35, GREEN, True, GREEN)
add_text(slide, "Download when needed", 9.72, 5.32, 2.5, 0.2, 10.5, WHITE, True, PP_ALIGN.CENTER)
add_footer(slide, 9)


# 10 — Control and trust
slide = prs.slides.add_slide(blank)
add_title(slide, "Control, trust and continuity", "Digital records improve both day-to-day discipline and long-term reliability.", 10)
controls = [
    ("Role-based access", "Desk, gate, coach and manager see only the pages assigned to them.", BLUE, "A"),
    ("Activity log", "Important changes can be traced to a user instead of relying on memory.", GREEN, "L"),
    ("Pool data isolation", "Each pool account has separate swimmers, payments and settings.", ORANGE, "D"),
    ("Consistent records", "Pass, invoice, payment and attendance use the same swimmer record.", CYAN, "1"),
    ("Cloud access", "Authorised staff can work through a browser without one physical register.", BLUE, "C"),
    ("Downloadable reports", "Lists and records can be exported for follow-up or offline review.", GREEN, "↓"),
]
for idx, (title, body, color, glyph) in enumerate(controls):
    col, row = idx % 3, idx // 3
    add_icon_card(slide, title, body, 0.62 + col * 4.17, 1.48 + row * 2.35, 3.82, 1.9, color, glyph)
add_footer(slide, 10)


# 11 — 30-day rollout
slide = prs.slides.add_slide(blank)
add_title(slide, "A practical 30-day transition", "Move gradually — old paper can remain an archive while new activity starts digitally.", 11)
weeks = [
    ("WEEK 1", "Set up", "Pool details, batches, pass types and staff access"),
    ("WEEK 2", "Start registration", "Use SwimIT for new swimmers and renewals"),
    ("WEEK 3", "Scan at gate", "Use QR pass and attendance for daily entry"),
    ("WEEK 4", "Review money", "Enter expenses and review the first balance sheet"),
]
for idx, (week, title, body) in enumerate(weeks):
    x = 0.68 + idx * 3.13
    color = [BLUE, CYAN, GREEN, ORANGE][idx]
    add_rect(slide, x, 1.72, 2.7, 3.55, WHITE, True, LINE)
    add_rect(slide, x + 0.78, 1.4, 1.14, 0.58, color, True, color)
    add_text(slide, week, x + 0.78, 1.56, 1.14, 0.21, 10.5, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, str(idx + 1), x + 0.91, 2.25, 0.88, 0.75, 31, color, True, PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.25, 3.18, 2.2, 0.4, 17, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, body, x + 0.3, 3.85, 2.1, 0.95, 12.5, MUTED, False, PP_ALIGN.CENTER)
    if idx < 3:
        add_text(slide, "→", x + 2.72, 2.8, 0.42, 0.32, 20, LINE, True, PP_ALIGN.CENTER)
add_rect(slide, 1.28, 5.85, 10.77, 0.62, PALE_ORANGE, True, PALE_ORANGE)
add_text(
    slide,
    "No need to type ten years of old registers. Begin with today’s new joins and renewals.",
    1.5,
    6.02,
    10.33,
    0.27,
    13.5,
    NAVY,
    True,
    PP_ALIGN.CENTER,
)
add_footer(slide, 11)


# 12 — Demo plan
slide = prs.slides.add_slide(blank)
add_title(slide, "The 5-minute live demonstration", "Show the complete value chain instead of explaining every menu.", 12)
steps = [
    ("1", "Register", "Create a swimmer record"),
    ("2", "Collect", "Take cash / UPI payment"),
    ("3", "Deliver", "WhatsApp pass + invoice"),
    ("4", "Scan", "Mark gate attendance"),
    ("5", "Review", "See credit in Balance Sheet"),
]
for idx, (num, title, body) in enumerate(steps):
    x = 0.58 + idx * 2.54
    color = [BLUE, CYAN, GREEN, ORANGE, NAVY][idx]
    add_rect(slide, x, 1.65, 2.2, 3.52, WHITE, True, LINE)
    add_rect(slide, x + 0.72, 1.35, 0.76, 0.76, color, True, color)
    add_text(slide, num, x + 0.72, 1.52, 0.76, 0.32, 18, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.2, 2.48, 1.8, 0.45, 17, NAVY, True, PP_ALIGN.CENTER)
    add_text(slide, body, x + 0.28, 3.2, 1.64, 0.8, 12.5, MUTED, False, PP_ALIGN.CENTER)
    if idx < 4:
        add_text(slide, "→", x + 2.18, 3.0, 0.35, 0.3, 18, CYAN, True, PP_ALIGN.CENTER)
add_rect(slide, 1.75, 5.76, 9.82, 0.72, PALE_GREEN, True, PALE_GREEN)
add_text(
    slide,
    "One swimmer. One payment. One scan. One financial record.",
    2.02,
    5.96,
    9.28,
    0.3,
    18,
    NAVY,
    True,
    PP_ALIGN.CENTER,
)
add_footer(slide, 12)


# 13 — Closing
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, 13.333, 7.5, NAVY)
slide.shapes.add_picture(str(PUBLIC / "swimit-wordmark.png"), Inches(4.65), Inches(0.62), width=Inches(4.03))
add_text(slide, "Paper keeps you busy.\nSwimIT keeps you in control.", 1.55, 2.23, 10.23, 1.2, 30, WHITE, True, PP_ALIGN.CENTER)
add_text(
    slide,
    "Registration · Pass · Invoice · Gate · Attendance · Expenses · Balance Sheet",
    1.5,
    3.83,
    10.33,
    0.42,
    15,
    RGBColor(188, 225, 247),
    False,
    PP_ALIGN.CENTER,
)
add_rect(slide, 4.35, 4.78, 4.63, 0.68, CYAN, True, CYAN)
add_text(slide, "Move your pool forward with SwimIT", 4.55, 4.96, 4.23, 0.27, 14, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, "Thank you", 5.6, 6.56, 2.13, 0.36, 15, WHITE, True, PP_ALIGN.CENTER)


prs.save(OUTPUT)
print(f"Created {OUTPUT}")
print(f"Slides: {len(prs.slides)}")
